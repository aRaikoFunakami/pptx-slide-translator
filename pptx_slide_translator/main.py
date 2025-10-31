import os
import sys
import asyncio
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.prompts import PromptTemplate
from langchain_openai import ChatOpenAI
from pydantic import BaseModel, Field
from typing import List
from tqdm.asyncio import tqdm

# OPENAI_API_KEYの存在チェック
def check_openai_api_key():
    if not os.getenv("OPENAI_API_KEY"):
        print("エラー: OPENAI_API_KEYが設定されていません。")
        print("環境変数OPENAI_API_KEYを設定してから実行してください。")
        sys.exit(1)

class TranslationItem(BaseModel):
    id: str = Field(description="text id")
    translated: str = Field(description="translated text")

class TranslationResult(BaseModel):
    translations: List[TranslationItem] = Field(
        description="list of translated items"
    )

def get_model_config():
    """
    現在の環境変数からモデル設定を取得する
    """
    model = os.getenv("OPENAI_MODEL", "gpt-4.1-mini")
    
    # OSSモデルかどうかを判定してBASEURLを設定
    if "oss" in model.lower() or "local" in model.lower() or "gemma" in model.lower():
        baseurl = os.getenv("OPENAI_BASEURL", "http://localhost:11434/v1")
    else:
        baseurl = os.getenv("OPENAI_BASEURL", "")
    
    print(f"Using OPENAI_MODEL: {model}")
    if baseurl:
        print(f"Using OPENAI_BASEURL: {baseurl}")
    
    return model, baseurl

async def translate_texts_openai_async(texts: List[str], target_lang="en") -> List[str]:
    """
    OpenAI APIでテキストリストを並列バッチ翻訳する（非同期版）
    texts: 翻訳対象テキストのリスト
    target_lang: 'en' or 'ja'
    戻り値: 翻訳文リスト（元textsと同じ順）
    
    10文字列ごとにバッチ処理し、各バッチを並列で実行します。
    進捗状況はtqdmで表示されます。
    """
    if not texts:
        return []
    
    # 空のテキストを除外してインデックスを記録
    non_empty_texts = []
    text_indices = []
    for i, text in enumerate(texts):
        if text.strip():
            non_empty_texts.append({
                "id": f"text_{i}",
                "text": text.strip()
            })
            text_indices.append(i)
    
    if not non_empty_texts:
        return [""] * len(texts)

    parser = JsonOutputParser(pydantic_object=TranslationResult)
    
    # モデル設定を動的に取得
    openai_model, openai_baseurl = get_model_config()
    
    if openai_baseurl:
        model = ChatOpenAI(temperature=0, model=openai_model, openai_api_base=openai_baseurl)
    else:
        model = ChatOpenAI(temperature=0, model=openai_model)

    batch_size = 10
    all_translations = []

    async def process_chunk(chunk, pbar=None):
        """単一チャンクを処理する"""
        items_json = json.dumps(
            [{"id": item["id"], "text": item["text"]} for item in chunk], ensure_ascii=False
        )
        
        # プロンプトテンプレートを作成
        prompt_template = PromptTemplate(
            template=(
                "You are a professional translator. Please translate each text in the following list to "
                + ("English" if target_lang == "en" else "Japanese")
                + ". Return the result as a strict JSON object.\n"
                "The output format must be:\n"
                '{{\n  "translations": [\n    {{"id": "text_0", "translated": "Translated text 1"}},\n    {{"id": "text_1", "translated": "Translated text 2"}},\n    ...\n  ]\n}}\n'
                "Do not include any explanations or extra text. Only output the JSON object in the specified format.\n"
                "Here is the list of items to translate (as JSON):\n{items_json}\n"
                "{format_instructions}"
            ),
            input_variables=["items_json"],
            partial_variables={"format_instructions": parser.get_format_instructions()},
        )
        chain = prompt_template | model | parser
        
        try:
            result = await chain.ainvoke({"items_json": items_json})
            
            if not (isinstance(result, dict) and "translations" in result):
                if pbar:
                    pbar.update(len(chunk))
                return [item["text"] for item in chunk]  # エラー時は元のテキストを返す
            
            translations = result["translations"]
            id_to_trans = {item["id"]: item["translated"] for item in translations}
            chunk_translations = []
            
            # 各チャンクアイテムの翻訳結果を取得
            for item in chunk:
                translated = id_to_trans.get(item["id"], item["text"])
                if not translated.strip():
                    translated = item["text"]  # 空の場合は元のテキストを使用
                chunk_translations.append(translated)
                
            # 進捗バーを更新
            if pbar:
                pbar.update(len(chunk))
                
            return chunk_translations
        except Exception as e:
            print(f"翻訳処理でエラーが発生しました: {e}")
            if pbar:
                pbar.update(len(chunk))
            return [item["text"] for item in chunk]  # エラー時は元のテキストを返す

    # 進捗バーを初期化
    total_texts = len(non_empty_texts)
    with tqdm(total=total_texts, desc="翻訳中", unit="texts") as pbar:
        # バッチに分割
        tasks = [
            process_chunk(non_empty_texts[i:i + batch_size], pbar)
            for i in range(0, len(non_empty_texts), batch_size)
        ]
        # 並列実行
        results = await asyncio.gather(*tasks)
        
    # 結果を統合
    for chunk_result in results:
        all_translations.extend(chunk_result)
    
    # 結果を元の順序に戻す
    final_results = [""] * len(texts)
    for i, original_index in enumerate(text_indices):
        if i < len(all_translations):
            final_results[original_index] = all_translations[i]
        else:
            final_results[original_index] = texts[original_index]
    
    return final_results

def translate_text(src_text: str, target_lang="en") -> str:
    """
    単一のテキストを翻訳する同期関数
    """
    if not src_text.strip():
        return src_text
    
    # 非同期関数を同期実行
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    try:
        result = loop.run_until_complete(translate_texts_openai_async([src_text], target_lang))
        return result[0] if result else src_text
    finally:
        loop.close()

def collect_texts_from_shape(shape, texts_to_translate, text_objects):
    """
    シェイプから再帰的にテキストを収集する
    グループ化されたオブジェクトにも対応
    """
    # グループ化されたシェイプの場合は再帰的に処理
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            collect_texts_from_shape(child_shape, texts_to_translate, text_objects)
        return
    
    # テーブルの処理
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        try:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts_to_translate.append(cell.text)
                        text_objects.append(cell)
        except Exception as e:
            print(f"テーブル処理でエラー: {e}")
        return
    
    # 通常のテキストフレームの処理
    if hasattr(shape, "text_frame") and shape.text_frame is not None:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    texts_to_translate.append(run.text)
                    text_objects.append(run)

async def translate_pptx_async(input_path: str, output_path: str, target_lang="en"):
    """
    PPTXファイルを翻訳する（非同期版・一括翻訳）
    """
    print(f"PPTXファイルを読み込み中: {input_path}")
    prs = Presentation(input_path)
    
    # 全てのテキストを収集（オブジェクトへの直接参照を保持）
    texts_to_translate = []
    text_objects = []  # 実際のテキストオブジェクトへの参照
    
    print("テキストを収集中...")
    for slide in prs.slides:
        # 各シェイプから再帰的にテキストを収集
        for shape in slide.shapes:
            collect_texts_from_shape(shape, texts_to_translate, text_objects)
    
    print(f"収集されたテキスト数: {len(texts_to_translate)} (全 {len(prs.slides)} スライド)")
    
    if not texts_to_translate:
        print("翻訳対象のテキストが見つかりませんでした。")
        prs.save(output_path)
        return
    
    # 一括翻訳実行
    print("翻訳を実行中...")
    translated_texts = await translate_texts_openai_async(texts_to_translate, target_lang)
    
    # 翻訳結果を元のオブジェクトに直接適用
    print("翻訳結果を適用中...")
    applied_count = 0
    for text_obj, translated_text in zip(text_objects, translated_texts):
        if translated_text.strip():
            try:
                text_obj.text = translated_text
                applied_count += 1
            except Exception as e:
                print(f"翻訳結果の適用でエラー: {e}")
    
    print(f"{applied_count}/{len(text_objects)} 件のテキストを翻訳しました")
    
    # 翻訳済みPPTXを保存
    print(f"翻訳済みファイルを保存中: {output_path}")
    prs.save(output_path)
    print("翻訳完了！")

def translate_pptx(input_path: str, output_path: str, target_lang="en"):
    """
    PPTXファイルを翻訳する（同期版）
    """
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    try:
        loop.run_until_complete(translate_pptx_async(input_path, output_path, target_lang))
    finally:
        loop.close()

def main():
    """コマンドラインエントリーポイント"""
    import argparse
    
    # OPENAI_API_KEYの存在チェック
    check_openai_api_key()
    
    parser = argparse.ArgumentParser(description="PPTX翻訳ツール")
    parser.add_argument("input_file", help="入力PPTXファイルのパス")
    parser.add_argument("-o", "--output", default="output_translated.pptx", 
                       help="出力PPTXファイルのパス (デフォルト: output_translated.pptx)")
    parser.add_argument("-l", "--lang", choices=["ja", "en"], default="en",
                       help="翻訳先言語 (ja: 日本語, en: 英語) (デフォルト: en)")
    parser.add_argument("-m", "--model", 
                       help="AIモデル名 (例: gpt-4.1-mini, gpt-oss:20b, local-llama)")
    parser.add_argument("-u", "--baseurl", 
                       help="OpenAI互換API のベースURL (例: http://localhost:11434/v1)")
    
    args = parser.parse_args()
    
    # AIモデルとBASEURLの設定
    if args.model:
        os.environ["OPENAI_MODEL"] = args.model
    
    if args.baseurl:
        os.environ["OPENAI_BASEURL"] = args.baseurl
    elif args.model and ("oss" in args.model.lower() or "local" in args.model.lower() or "gemma" in args.model.lower()):
        # OSSモデルでBASEURLが指定されていない場合のデフォルト値
        if not os.getenv("OPENAI_BASEURL"):
            os.environ["OPENAI_BASEURL"] = "http://localhost:11434/v1"
    
    if not os.path.exists(args.input_file):
        print(f"エラー: 入力ファイルが見つかりません: {args.input_file}")
        sys.exit(1)
    
    print(f"入力ファイル: {args.input_file}")
    print(f"出力ファイル: {args.output}")
    print(f"翻訳先言語: {'日本語' if args.lang == 'ja' else '英語'}")
    print("-" * 50)
    
    translate_pptx(args.input_file, args.output, args.lang)

if __name__ == "__main__":
    main()
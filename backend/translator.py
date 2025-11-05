"""
PPTXファイル翻訳ロジック
既存のmain.pyから翻訳機能を分離
"""
import asyncio
import json
import os
from typing import List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.prompts import PromptTemplate
from langchain_openai import ChatOpenAI
from pydantic import BaseModel, Field


class TranslationItem(BaseModel):
    id: str = Field(description="text id")
    translated: str = Field(description="translated text")


class TranslationResult(BaseModel):
    translations: List[TranslationItem] = Field(
        description="list of translated items"
    )


def get_model_config():
    """現在の環境変数からモデル設定を取得する"""
    model = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
    baseurl = os.getenv("OPENAI_BASEURL", "")
    
    return model, baseurl


async def translate_texts_openai_async(texts: List[str], target_lang="en") -> List[str]:
    """
    OpenAI APIでテキストリストを並列バッチ翻訳する（非同期版）
    """
    if not texts:
        return []
    
    # 空のテキストを除外してインデックスを記録
    non_empty_texts = []
    text_indices = []
    for i, text in enumerate(texts):
        if text.strip():
            non_empty_texts.append({
                "id": f"text_{i}",
                "text": text.strip()
            })
            text_indices.append(i)
    
    if not non_empty_texts:
        return [""] * len(texts)

    parser = JsonOutputParser(pydantic_object=TranslationResult)
    
    # モデル設定を動的に取得
    openai_model, openai_baseurl = get_model_config()
    
    if openai_baseurl:
        model = ChatOpenAI(temperature=0, model=openai_model, openai_api_base=openai_baseurl)
    else:
        model = ChatOpenAI(temperature=0, model=openai_model)

    batch_size = 10
    all_translations = []

    async def process_chunk(chunk):
        """単一チャンクを処理する"""
        items_json = json.dumps(
            [{"id": item["id"], "text": item["text"]} for item in chunk], ensure_ascii=False
        )
        
        # プロンプトテンプレートを作成
        prompt_template = PromptTemplate(
            template=(
                "You are a professional translator. Please translate each text in the following list to "
                + ("English" if target_lang == "en" else "Japanese")
                + ". Return the result as a strict JSON object.\n"
                "The output format must be:\n"
                '{{\n  "translations": [\n    {{"id": "text_0", "translated": "Translated text 1"}},\n    {{"id": "text_1", "translated": "Translated text 2"}},\n    ...\n  ]\n}}\n'
                "Do not include any explanations or extra text. Only output the JSON object in the specified format.\n"
                "Here is the list of items to translate (as JSON):\n{items_json}\n"
                "{format_instructions}"
            ),
            input_variables=["items_json"],
            partial_variables={"format_instructions": parser.get_format_instructions()},
        )
        chain = prompt_template | model | parser
        
        try:
            result = await chain.ainvoke({"items_json": items_json})
            
            if not (isinstance(result, dict) and "translations" in result):
                return [item["text"] for item in chunk]  # エラー時は元のテキストを返す
            
            translations = result["translations"]
            id_to_trans = {item["id"]: item["translated"] for item in translations}
            chunk_translations = []
            
            # 各チャンクアイテムの翻訳結果を取得
            for item in chunk:
                translated = id_to_trans.get(item["id"], item["text"])
                if not translated.strip():
                    translated = item["text"]  # 空の場合は元のテキストを使用
                chunk_translations.append(translated)
                
            return chunk_translations
        except Exception as e:
            print(f"翻訳処理でエラーが発生しました: {e}")
            return [item["text"] for item in chunk]  # エラー時は元のテキストを返す

    # バッチに分割して並列実行
    tasks = [
        process_chunk(non_empty_texts[i:i + batch_size])
        for i in range(0, len(non_empty_texts), batch_size)
    ]
    results = await asyncio.gather(*tasks)
    
    # 結果を統合
    for chunk_result in results:
        all_translations.extend(chunk_result)
    
    # 結果を元の順序に戻す
    final_results = [""] * len(texts)
    for i, original_index in enumerate(text_indices):
        if i < len(all_translations):
            final_results[original_index] = all_translations[i]
        else:
            final_results[original_index] = texts[original_index]
    
    return final_results


def collect_texts_from_shape(shape, texts_to_translate, text_objects):
    """
    シェイプから再帰的にテキストを収集する
    グループ化されたオブジェクトにも対応
    """
    # グループ化されたシェイプの場合は再帰的に処理
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            collect_texts_from_shape(child_shape, texts_to_translate, text_objects)
        return
    
    # テーブルの処理
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        try:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts_to_translate.append(cell.text)
                        text_objects.append(cell)
        except Exception as e:
            print(f"テーブル処理でエラー: {e}")
        return
    
    # 通常のテキストフレームの処理
    if hasattr(shape, "text_frame") and shape.text_frame is not None:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    texts_to_translate.append(run.text)
                    text_objects.append(run)


async def translate_pptx_async(input_path: str, output_path: str, target_lang="en"):
    """
    PPTXファイルを翻訳する（非同期版・一括翻訳）
    """
    prs = Presentation(input_path)
    
    # 全てのテキストを収集（オブジェクトへの直接参照を保持）
    texts_to_translate = []
    text_objects = []  # 実際のテキストオブジェクトへの参照
    
    for slide in prs.slides:
        # 各シェイプから再帰的にテキストを収集
        for shape in slide.shapes:
            collect_texts_from_shape(shape, texts_to_translate, text_objects)
    
    if not texts_to_translate:
        prs.save(output_path)
        return len(prs.slides), 0
    
    # 一括翻訳実行
    translated_texts = await translate_texts_openai_async(texts_to_translate, target_lang)
    
    # 翻訳結果を元のオブジェクトに直接適用
    applied_count = 0
    for text_obj, translated_text in zip(text_objects, translated_texts):
        if translated_text.strip():
            try:
                text_obj.text = translated_text
                applied_count += 1
            except Exception as e:
                print(f"翻訳結果の適用でエラー: {e}")
    
    # 翻訳済みPPTXを保存
    prs.save(output_path)
    
    return len(prs.slides), applied_count


def analyze_pptx(file_path: str) -> dict:
    """
    PPTXファイルを分析してページ数と翻訳対象テキスト数を返す
    """
    try:
        prs = Presentation(file_path)
        
        texts_to_translate = []
        text_objects = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                collect_texts_from_shape(shape, texts_to_translate, text_objects)
        
        return {
            "pages": len(prs.slides),
            "text_count": len(texts_to_translate),
            "success": True
        }
    except Exception as e:
        return {
            "pages": 0,
            "text_count": 0,
            "success": False,
            "error": str(e)
        }